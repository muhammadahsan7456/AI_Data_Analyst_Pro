import os
import io
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors


def export_to_csv(df: pd.DataFrame) -> bytes:
    """
    Export DataFrame to CSV bytes stream.
    """
    return df.to_csv(index=False).encode('utf-8')


def export_to_excel(df: pd.DataFrame, dataset_name: str = "Dataset") -> bytes:
    """
    Export DataFrame to styled Excel (.xlsx) bytes stream.
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    header_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    cell_alignment = Alignment(horizontal="left", vertical="center")
    thin_border = Border(
        left=Side(style='thin', color='E5E7EB'),
        right=Side(style='thin', color='E5E7EB'),
        top=Side(style='thin', color='E5E7EB'),
        bottom=Side(style='thin', color='E5E7EB')
    )

    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
        ws.append(row)
        for c_idx in range(1, len(row) + 1):
            cell = ws.cell(row=r_idx, column=c_idx)
            cell.alignment = cell_alignment
            cell.border = thin_border
            if r_idx == 1:
                cell.fill = header_fill
                cell.font = header_font

    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = col[0].column_letter
        ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

    ws_summary = wb.create_sheet(title="Summary")
    ws_summary.append(["Property", "Value"])
    ws_summary.cell(1, 1).font = header_font
    ws_summary.cell(1, 2).font = header_font

    ws_summary.append(["Dataset Name", dataset_name])
    ws_summary.append(["Total Records", len(df)])
    ws_summary.append(["Total Columns", len(df.columns)])

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def export_to_pdf_report(df: pd.DataFrame, dataset_name: str = "Dataset", insights: list = None, chart_file_path: str = None) -> bytes:
    """
    Generate Executive PDF report with Logo, KPIs, Insights, Charts, and ALL Query Records (Multi-Page Paginated Table).
    """
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'TitleCustom',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        leading=24,
        textColor=colors.HexColor('#1E3A8A'),
        spaceAfter=10
    )

    h2_style = ParagraphStyle(
        'Heading2Custom',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor('#2563EB'),
        spaceBefore=12,
        spaceAfter=6
    )

    body_style = ParagraphStyle(
        'BodyCustom',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12,
        textColor=colors.HexColor('#374151'),
        spaceAfter=4
    )

    table_hdr_style = ParagraphStyle(
        'TableHdr',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8,
        leading=10,
        textColor=colors.HexColor('#1F2937')
    )

    table_cell_style = ParagraphStyle(
        'TableCell',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=8,
        leading=10,
        textColor=colors.HexColor('#4B5563')
    )

    elements = []

    elements.append(Paragraph("Executive Data Analysis Report", title_style))
    elements.append(Paragraph(f"<b>Dataset / Query:</b> {dataset_name} | <b>Total Records:</b> {len(df)} | <b>Total Columns:</b> {len(df.columns)}", body_style))
    elements.append(Spacer(1, 8))

    if insights:
        elements.append(Paragraph("🤖 Key AI Insights & Business Recommendations", h2_style))
        for insight in insights:
            elements.append(Paragraph(f"• {insight}", body_style))
        elements.append(Spacer(1, 8))

    if chart_file_path:
        full_chart_path = os.path.join("static", chart_file_path) if not chart_file_path.startswith("static") else chart_file_path
        if os.path.exists(full_chart_path):
            elements.append(Paragraph("📈 Visual Analysis", h2_style))
            img = RLImage(full_chart_path, width=500, height=220)
            elements.append(img)
            elements.append(Spacer(1, 10))

    elements.append(Paragraph(f"📋 Complete Records View ({len(df)} Matching Rows)", h2_style))

    # Process ALL rows of df (no truncation) - Fast tuple iteration
    cols_to_render = df.columns[:8]
    col_indices = [df.columns.get_loc(c) for c in cols_to_render]
    col_width = 540 / max(len(cols_to_render), 1)

    table_data = [[Paragraph(f"<b>{col}</b>", table_hdr_style) for col in cols_to_render]]
    for row_tuple in df.itertuples(index=False):
        row_cells = []
        for idx in col_indices:
            val_raw = row_tuple[idx]
            val = str(val_raw) if val_raw is not None and not pd.isna(val_raw) else ""
            if len(val) > 40:
                val = val[:37] + "..."
            row_cells.append(Paragraph(val, table_cell_style))
        table_data.append(row_cells)

    pdf_table = Table(
        table_data,
        colWidths=[col_width] * len(cols_to_render),
        repeatRows=1,
        hAlign='LEFT'
    )
    pdf_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F3F4F6')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#1F2937')),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
        ('TOPPADDING', (0, 0), (-1, -1), 3),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#E5E7EB')),
    ]))
    elements.append(pdf_table)

    doc.build(elements)
    buffer.seek(0)
    return buffer.getvalue()


def export_to_word_report(df: pd.DataFrame, dataset_name: str = "Dataset", insights: list = None, chart_file_path: str = None) -> bytes:
    """
    Export Executive Word (.docx) report with AI insights, chart, and sample table.
    """
    try:
        import docx
        from docx.shared import Inches
        doc = docx.Document()
        doc.add_heading(f"Executive Data Report: {dataset_name}", 0)
        doc.add_paragraph(f"Total Records: {len(df)} | Total Columns: {len(df.columns)}")

        if insights:
            doc.add_heading("AI Business Insights & Recommendations", level=1)
            for ins in insights:
                doc.add_paragraph(f"• {ins}")

        if chart_file_path:
            full_chart_path = os.path.join("static", chart_file_path) if not chart_file_path.startswith("static") else chart_file_path
            if os.path.exists(full_chart_path):
                doc.add_heading("Visual Analysis Chart", level=1)
                doc.add_picture(full_chart_path, width=Inches(5.5))

        doc.add_heading("Sample Data Preview (Top 10 Rows)", level=1)
        table = doc.add_table(rows=1, cols=min(len(df.columns), 6))
        hdr_cells = table.rows[0].cells
        for i, col in enumerate(df.columns[:6]):
            hdr_cells[i].text = str(col)

        for _, row in df.head(10).iterrows():
            row_cells = table.add_row().cells
            for i, col in enumerate(df.columns[:6]):
                row_cells[i].text = str(row[col])[:30]

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    except Exception:
        html_doc = f"""
        <html>
        <head><meta charset="utf-8"><title>{dataset_name} Executive Report</title></head>
        <body style="font-family: Arial, sans-serif; padding: 20px;">
            <h1 style="color: #1e3a8a;">Executive Data Report: {dataset_name}</h1>
            <p><strong>Total Records:</strong> {len(df)} | <strong>Total Columns:</strong> {len(df.columns)}</p>
            {"<h2>AI Insights</h2><ul>" + "".join(f"<li>{i}</li>" for i in insights) + "</ul>" if insights else ""}
            <h2>Sample Dataset</h2>
            {df.head(15).to_html(classes="table", index=False)}
        </body>
        </html>
        """
        return html_doc.encode("utf-8")


def export_to_pptx_report(df: pd.DataFrame, dataset_name: str = "Dataset", insights: list = None, chart_file_path: str = None) -> bytes:
    """
    Export Executive PowerPoint (.pptx) presentation with Title, Insights, Embedded Chart, and Formatted Dataset Table Slides.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # Slide 1: Title Slide
        slide1 = prs.slides.add_slide(blank_layout)
        bg1 = slide1.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = RGBColor(15, 23, 42)
        bg1.line.color.rgb = RGBColor(15, 23, 42)

        txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True

        p1 = tf1.paragraphs[0]
        p1.text = "EXECUTIVE DATA ANALYSIS REPORT"
        p1.font.bold = True
        p1.font.size = Pt(36)
        p1.font.color.rgb = RGBColor(56, 189, 248)

        p2 = tf1.add_paragraph()
        p2.text = f"Dataset: {dataset_name}"
        p2.font.bold = True
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(248, 250, 252)

        p3 = tf1.add_paragraph()
        p3.text = f"Total Records: {len(df):,} | Total Columns: {len(df.columns)} | AI Data Analyst Pro Enterprise Platform"
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(148, 163, 184)

        # Slide 2: AI Insights Slide
        if insights:
            slide2 = prs.slides.add_slide(blank_layout)
            bg2 = slide2.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
            bg2.fill.solid()
            bg2.fill.fore_color.rgb = RGBColor(15, 23, 42)
            bg2.line.color.rgb = RGBColor(15, 23, 42)

            tx_hdr2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
            p_hdr2 = tx_hdr2.text_frame.paragraphs[0]
            p_hdr2.text = "🤖 AI Key Business Insights & Recommendations"
            p_hdr2.font.bold = True
            p_hdr2.font.size = Pt(24)
            p_hdr2.font.color.rgb = RGBColor(56, 189, 248)

            tx_ins2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5))
            tf_ins2 = tx_ins2.text_frame
            tf_ins2.word_wrap = True
            for ins in insights[:6]:
                p = tf_ins2.add_paragraph()
                p.text = f"• {ins}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(226, 232, 240)
                p.space_after = Pt(12)

        # Slide 3: Chart Slide
        if chart_file_path:
            full_chart_path = os.path.join("static", chart_file_path) if not chart_file_path.startswith("static") else chart_file_path
            if os.path.exists(full_chart_path):
                slide_c = prs.slides.add_slide(blank_layout)
                bg_c = slide_c.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
                bg_c.fill.solid()
                bg_c.fill.fore_color.rgb = RGBColor(15, 23, 42)
                bg_c.line.color.rgb = RGBColor(15, 23, 42)

                tx_c = slide_c.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
                p_c = tx_c.text_frame.paragraphs[0]
                p_c.text = "📈 Executive Visual Analytics Chart"
                p_c.font.bold = True
                p_c.font.size = Pt(24)
                p_c.font.color.rgb = RGBColor(56, 189, 248)

                slide_c.shapes.add_picture(full_chart_path, Inches(1.5), Inches(1.5), width=Inches(10.333))

        # Slide 4: Complete Dataset Records View Table Slide
        if not df.empty:
            slide3 = prs.slides.add_slide(blank_layout)
            bg3 = slide3.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
            bg3.fill.solid()
            bg3.fill.fore_color.rgb = RGBColor(15, 23, 42)
            bg3.line.color.rgb = RGBColor(15, 23, 42)

            tx_t3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
            p_t3 = tx_t3.text_frame.paragraphs[0]
            p_t3.text = f"📋 Dataset Records View ({len(df)} Records)"
            p_t3.font.bold = True
            p_t3.font.size = Pt(24)
            p_t3.font.color.rgb = RGBColor(56, 189, 248)

            cols_to_show = list(df.columns[:7])
            rows_to_show = df.head(15)

            num_rows = len(rows_to_show) + 1
            num_cols = len(cols_to_show)

            table_shape3 = slide3.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
            table3 = table_shape3.table

            for c_idx, col_name in enumerate(cols_to_show):
                cell = table3.cell(0, c_idx)
                cell.text = str(col_name)[:18]
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(255, 255, 255)

            for r_idx, (_, row) in enumerate(rows_to_show.iterrows(), start=1):
                for c_idx, col_name in enumerate(cols_to_show):
                    cell = table3.cell(r_idx, c_idx)
                    val = str(row[col_name]) if row[col_name] is not None and not pd.isna(row[col_name]) else ""
                    cell.text = val[:25]
                    cell.fill.solid()
                    if r_idx % 2 == 0:
                        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    else:
                        cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = RGBColor(226, 232, 240)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as pptx_err:
        print("[PPTX ENGINE FALLBACK]", pptx_err)
        html_ppt = f"""
        <html>
        <head><meta charset="utf-8"><title>{dataset_name} Executive Presentation</title></head>
        <body style="font-family: Arial, sans-serif; padding: 40px; background: #0f172a; color: white;">
            <div style="border: 2px solid #38bdf8; padding: 40px; border-radius: 12px; margin-bottom: 20px;">
                <h1 style="color: #38bdf8;">Data Analysis: {dataset_name}</h1>
                <p>Total Records: {len(df)} | Total Columns: {len(df.columns)}</p>
            </div>
            {"<div style='border: 2px solid #a855f7; padding: 40px; border-radius: 12px;'><h2>AI Insights</h2><ul>" + "".join(f"<li>{i}</li>" for i in insights[:5]) + "</ul></div>" if insights else ""}
            <div style="margin-top: 20px;">
                <h2>Dataset Sample Table</h2>
                {df.head(15).to_html(classes="table", index=False)}
            </div>
        </body>
        </html>
        """
        return html_ppt.encode("utf-8")
